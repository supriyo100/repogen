from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
HEADER_BG = RGBColor(0, 51, 102)  # Dark Blue
ACCENT = RGBColor(0, 102, 204)   # Bright Blue
TEXT = RGBColor(51, 51, 51)      # Dark Gray
SUCCESS = RGBColor(34, 139, 34)  # Green
WARNING = RGBColor(220, 20, 60)  # Red
LIGHT_BG = RGBColor(240, 248, 255)  # Alice Blue

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = HEADER_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)

def add_content_slide(prs, title, content_items):
    """Add content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = HEADER_BG
    title_shape.line.color.rgb = HEADER_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    y_pos = 1.0
    for item in content_items:
        if isinstance(item, dict):
            item_type = item.get('type', 'text')
            
            if item_type == 'heading':
                # Section heading
                heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.35))
                heading_frame = heading_box.text_frame
                p = heading_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = ACCENT
                y_pos += 0.35
                
            elif item_type == 'bullet':
                # Bullet point
                bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.7), Inches(0.25))
                bullet_frame = bullet_box.text_frame
                bullet_frame.word_wrap = True
                p = bullet_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT
                p.level = item.get('level', 0)
                y_pos += 0.25
                
            elif item_type == 'table_row':
                # Table row
                row_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(8.8), Inches(0.3))
                row_frame = row_box.text_frame
                row_frame.word_wrap = True
                p = row_frame.paragraphs[0]
                p.text = f"  {item['col1']:<30} → {item['col2']:<25} {item['col3']}"
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT
                p.font.name = 'Courier New'
                y_pos += 0.30

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = HEADER_BG

title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.paragraphs[0]
p.text = "🏭 Biocon Malaysia\nProduction Planning Model"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
p = subtitle_frame.paragraphs[0]
p.text = "Model Assumptions, Rules & Considerations"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(200, 220, 255)

# Slide 2: TAT & Batch Assumptions
add_content_slide(prs, "1️⃣ TAT & Batch Assumptions", [
    {'type': 'heading', 'text': '⏱️ Turnaround Time (TAT)'},
    {'type': 'bullet', 'text': 'Standard Batch Size: ~1,150 kg (~9 pallets)'},
    {'type': 'bullet', 'text': 'Processing Time: 2.5 days (8 hours per pallet)'},
    {'type': 'bullet', 'text': 'Current consideration: Batch granularity in scheduling'},
    {'type': 'bullet', 'text': '⚠ Exception: Aspart & Canada = 900 kg (TAT rules unavailable)'},
    {'type': 'heading', 'text': '📦 Batch Switching Rule'},
    {'type': 'bullet', 'text': 'Minimum 5 batches required before switching between vial/pen'},
    {'type': 'bullet', 'text': 'Model considers demand fulfillment respectively'},
    {'type': 'bullet', 'text': 'Objective: Reduce changeover costs & improve line utilization'},
])

# Slide 3: Routing Rules
add_content_slide(prs, "2️⃣ Routing for Packaging Resources", [
    {'type': 'heading', 'text': '📍 Pack Size → Routing Rules'},
    {'type': 'table_row', 'col1': 'Pack of 1/10 (Low Vol)', 'col2': 'Manual Pack', 'col3': '(MFMPPL012702001)'},
    {'type': 'table_row', 'col1': 'Pack of 5 (DLP)', 'col2': 'Assembly + Auto Pack', 'col3': '(MFAPL0012702001)'},
    {'type': 'table_row', 'col1': 'Aspart/Basalog (Pack 5)', 'col2': 'MARCHESINI Line', 'col3': '(MFAPL0012702001)'},
    {'type': 'table_row', 'col1': 'Vial/Cart (Auto)', 'col2': 'MARCHESINI (Reused)', 'col3': '(MFAPL0012702001)'},
    {'type': 'heading', 'text': '🏢 Resource Assignment'},
    {'type': 'table_row', 'col1': 'EU DLP', 'col2': 'Auto Pack EU', 'col3': '(MFAPPL012702001)'},
    {'type': 'table_row', 'col1': 'US/EM DLP', 'col2': 'Auto Pack IMA', 'col3': '(MFIPL0012702001)'},
    {'type': 'table_row', 'col1': 'Manual Low Vol', 'col2': 'Manual Pack Line', 'col3': '(MFMPPL012702001)'},
    {'type': 'table_row', 'col1': 'Inspection', 'col2': 'AVI/MVI Shared', 'col3': '(03-02-AVI-02)'},
])
# Slide 3: Routing Rules (UPDATED)
add_content_slide(prs, "2️⃣ Routing for Packaging Resources", [
    {'type': 'heading', 'text': '🎯 Reason for Routing Rule Creation'},
    {'type': 'bullet', 'text': '❌ Gap Identified: Not all pegged 8-series Market SKUs mapped to resources in "DP Line Utilization" (SNP.xlsx)'},
    {'type': 'bullet', 'text': '🔍 Methodology: Analyzed market, type, and quantity of SKUs pegged with resources'},
    {'type': 'bullet', 'text': '✅ Solution: Mapped similar SKUs to similar resources across production stages using pattern matching'},
    
    {'type': 'heading', 'text': '📍 Pack Size → Routing Rules'},
    {'type': 'table_row', 'col1': 'Pack of 1/10 (Low Vol)', 'col2': 'Manual Pack', 'col3': '(MFMPPL012702001)'},
    {'type': 'table_row', 'col1': 'Pack of 5 (DLP)', 'col2': 'Assembly + Auto Pack', 'col3': '(MFAPL0012702001)'},
    {'type': 'table_row', 'col1': 'Aspart/Basalog (Pack 5)', 'col2': 'MARCHESINI Line', 'col3': '(MFAPL0012702001)'},
    {'type': 'table_row', 'col1': 'Vial/Cart (Auto)', 'col2': 'MARCHESINI (Reused)', 'col3': '(MFAPL0012702001)'},
    
    {'type': 'heading', 'text': '🏢 Resource Assignment'},
    {'type': 'table_row', 'col1': 'EU DLP', 'col2': 'Auto Pack EU', 'col3': '(MFAPPL012702001)'},
    {'type': 'table_row', 'col1': 'US/EM DLP', 'col2': 'Auto Pack IMA', 'col3': '(MFIPL0012702001)'},
    {'type': 'table_row', 'col1': 'Manual Low Vol', 'col2': 'Manual Pack Line', 'col3': '(MFMPPL012702001)'},
    {'type': 'table_row', 'col1': 'Inspection', 'col2': 'AVI/MVI Shared', 'col3': '(03-02-AVI-02)'},
])
# Slide 4: SKU & Demand Mapping
add_content_slide(prs, "3️⃣ SKU Demand Mapping Challenge", [
    {'type': 'heading', 'text': '📊 Demand Coverage Status'},
    {'type': 'bullet', 'text': '✔ Only 13 market SKUs with demand available (L2Ph1_Detail)'},
    {'type': 'bullet', 'text': '❌ Not all demand-mapped SKUs have BOM in DP Shortage'},
    {'type': 'bullet', 'text': '⚠ Demand too low → resources appeared non-constrained & idle'},
    {'type': 'heading', 'text': '🔧 Mitigation Approach'},
    {'type': 'bullet', 'text': '✔ Added 4 more market SKUs from ADV & EM Mar 2025 sheets'},
    {'type': 'bullet', 'text': '✔ Added demand for additional SKUs (no BOM initially)'},
    {'type': 'bullet', 'text': '🎯 Target: Scale demand to 2–3M EA/month (realistic load)'},
    {'type': 'bullet', 'text': '📈 Goal: Model behaves like Biocon real plant utilization'},
])

# Slide 5: Synthetic BOM Rules
add_content_slide(prs, "4️⃣ Synthetic BOM Creation Rules", [
    {'type': 'heading', 'text': '🔗 Common BOM Inheritance (Glargine DLP EU)'},
    {'type': 'bullet', 'text': 'SKUs in same family share large BOM portion'},
    {'type': 'bullet', 'text': 'Common materials: shipper boxes, COOL INSUL, 3-ply board, foil, BOPP film'},
    {'type': 'bullet', 'text': 'Example: 800006526 & 800006527 share 650003790, 650005260, etc.'},
    {'type': 'heading', 'text': '📝 SKU-Specific Variations'},
    {'type': 'bullet', 'text': 'Printed materials: Pen labels (379.5K), Cartons (75.9K), PIL, IFU'},
    {'type': 'bullet', 'text': 'Market-specific codes (Germany, Nordics, Austria)'},
    {'type': 'bullet', 'text': 'Blank labels vary by SKU: 1,130–2,500 per batch'},
    {'type': 'heading', 'text': '✨ Inheritance for Unmapped SKUs'},
    {'type': 'bullet', 'text': 'New markets inherit: Common BOM + structure + scaled quantities'},
    {'type': 'bullet', 'text': 'Shipper labels: Use nearest geography pattern'},
])

# Slide 6: Batch Scaling & Purpose
add_content_slide(prs, "5️⃣ Batch Scaling & Synthetic BOM Purpose", [
    {'type': 'heading', 'text': '📏 Batch Scaling Logic'},
    {'type': 'bullet', 'text': 'Reference: 330,000-unit batch (standard)'},
    {'type': 'bullet', 'text': 'Scaling: Linear scaling except pallet/shipper-based items'},
    {'type': 'bullet', 'text': 'Impact: Allows flexible batch sizing in model'},
    {'type': 'heading', 'text': '🎯 Why Synthetic BOM?'},
    {'type': 'bullet', 'text': 'Expand downstream pegging & supply chain visibility'},
    {'type': 'bullet', 'text': 'Increase SKU count for realistic planning runs'},
    {'type': 'bullet', 'text': 'Simulate real-world complexity (demand ≠ BOM coverage)'},
    {'type': 'bullet', 'text': 'Test system behavior under realistic demand loads'},
    {'type': 'bullet', 'text': 'Enable accurate resource utilization modeling'},
])

# Slide 7: Batch Mapping & Current Status
add_content_slide(prs, "6️⃣ Batch Mapping Clarification", [
    {'type': 'heading', 'text': '🔄 7-Series to 8-Series Mapping'},
    {'type': 'bullet', 'text': '❌ NOT 1:1 mapping (as initially assumed)'},
    {'type': 'bullet', 'text': 'Reality: One 7-series SFG batch splits into multiple 8-series FG SKUs'},
    {'type': 'bullet', 'text': 'Status: True allocation logic not yet provided by Biocon'},
    {'type': 'heading', 'text': '⏸️ Current Interim Approach'},
    {'type': 'bullet', 'text': 'Continue with 1:1 assumption for demo run'},
    {'type': 'bullet', 'text': 'Incorporate real mapping rules in Phase 2'},
    {'type': 'bullet', 'text': '⚠ Awaiting clarification from Biocon on actual split rules'},
])

# Slide 8: Key Considerations
add_content_slide(prs, "7️⃣ Model Considerations & Limitations", [
    {'type': 'heading', 'text': '✅ Implemented'},
    {'type': 'bullet', 'text': 'Daily scheduling with batch granularity'},
    {'type': 'bullet', 'text': 'Resource routing & line assignments'},
    {'type': 'bullet', 'text': 'Changeover rules (molecule 24h, type 6h)'},
    {'type': 'bullet', 'text': 'Demand fulfillment constraints'},
    {'type': 'heading', 'text': '⏳ Pending'},
    {'type': 'bullet', 'text': 'Accurate 7-series → 8-series split logic'},
    {'type': 'bullet', 'text': 'Complete BOM for all demand-mapped SKUs'},
    {'type': 'bullet', 'text': 'Real-time demand data from Q4 2025'},
    {'type': 'bullet', 'text': 'Validation of synthetic BOM assumptions'},
])

# Slide 9: Summary & Next Steps
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide9.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BG

title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "8️⃣ Summary & Next Steps"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = HEADER_BG

content_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

items = [
    ("✅ ASSUMPTIONS FINALIZED", [
        "TAT: 1,150 kg / 2.5 days",
        "Min 5 batch rule before format switch",
        "Routing rules for all pack sizes",
        "Synthetic BOM creation logic"
    ]),
    ("🔧 IN PROGRESS", [
        "Demand scaling to 2–3M EA/month",
        "BOM coverage expansion",
        "Resource utilization analysis"
    ]),
    ("⏳ NEXT PHASE", [
        "Confirm 7→8 series split logic",
        "Validate synthetic BOM vs actual",
        "Full Q4 2025 demand integration",
        "Model refinement & validation run"
    ])
]

for section_title, section_items in items:
    p = content_frame.add_paragraph()
    p.text = section_title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.level = 0
    
    for item in section_items:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.level = 1

# Save
output_path = "/home/supriyo/Downloads/Biocon_Model_Assumptions_Rules_Considerations.pptx"
prs.save(output_path)
print(f"✅ PowerPoint created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📌 Content:")
print(f"   1. Title Slide")
print(f"   2. TAT & Batch Assumptions")
print(f"   3. Routing Rules for Packaging")
print(f"   4. SKU Demand Mapping")
print(f"   5. Synthetic BOM Rules")
print(f"   6. Batch Scaling & Purpose")
print(f"   7. Batch Mapping Clarification")
print(f"   8. Model Considerations")
print(f"   9. Summary & Next Steps")