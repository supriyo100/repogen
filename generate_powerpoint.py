from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
HEADER_COLOR = RGBColor(0, 51, 102)  # Dark Blue
ACCENT_COLOR = RGBColor(0, 102, 204)  # Bright Blue
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark Gray
SUCCESS_COLOR = RGBColor(34, 139, 34)  # Forest Green
WARNING_COLOR = RGBColor(220, 20, 60)  # Crimson

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = HEADER_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Status Update: {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(prs, title, sections):
    """Add content slide with sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = HEADER_COLOR
    
    # Content
    y_pos = 1.1
    for section_title, items, icon in sections:
        # Section header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
        header_frame = header_box.text_frame
        p = header_frame.paragraphs[0]
        p.text = f"{icon} {section_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        y_pos += 0.35
        
        # Items
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.7), Inches(0.28))
            item_frame = item_box.text_frame
            item_frame.word_wrap = True
            p = item_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_COLOR
            
            y_pos += 0.28
        
        y_pos += 0.15

# Slide 1: Title
add_title_slide(prs, 
    "📘 Biocon Malaysia", 
    "Daily Production Planning Model – Status Update")

# Slide 2: Current Status
add_content_slide(prs, "1️⃣ Current Status – Where We Stand", [
    ("Material-BOM Pegging", [
        "✔ Material-BOM Pegging shared",
        "✔ Master Process & Planning Assumptions delivered",
        "✔ Biocon clarified: Batch mapping 7-series → 8-series = 1:1",
        "✔ Changeover times: Molecule 24h, Type 6h, Both = 24h (not cumulative)"
    ], "✅"),
    ("ORMAE Progress", [
        "✔ Generated input-output templates",
        "✔ Completed initial model test runs",
        "✔ Built PowerBI visualization for outputs"
    ], "🚀")
])

# Slide 3: Initial Observations
add_content_slide(prs, "2️⃣ Initial Model Observations", [
    ("Model Status", [
        "✔ Model structurally working (daily scheduling, batch allocation, resource assignment)",
        "⚠ Resource capacity not binding → demand mapped is still too low",
        "🔧 Need full Oct–Dec 2025 demand (~2–3M units/month) reflected accurately"
    ], "📊")
])

# Slide 4: Key Blockers
add_content_slide(prs, "3️⃣ Key Pending / Blockers", [
    ("Clarifications Needed", [
        "❗ Changeover clarification: For filling SKUs with same molecule & format",
        "   (e.g., 700001012 → 700004130), Is a changeover required?",
        "   → Awaiting Nisalini's confirmation",
        "❗ Incomplete demand/BOM for certain SKUs",
        "   → Synthetic BOM assumptions required"
    ], "⚠️")
])

# Slide 5: Next Steps
add_content_slide(prs, "4️⃣ Next Steps", [
    ("Action Items", [
        "🔄 Map complete Q4 2025 demand across all 8-series SKUs",
        "🧩 Add synthetic BOMs for undefined SKUs (per agreed rules)",
        "🚀 Re-run model with full realistic demand load",
        "📊 Review updated results with Iris on Monday, incorporate feedback",
        "👥 ORMAE to present consolidated run in alignment meeting with Mahesh"
    ], "📋")
])

# Slide 6: Collaboration & Concerns
add_content_slide(prs, "5️⃣ Collaboration & Concerns", [
    ("Team Alignment", [
        "🤝 Collaboration with ORMAE: smooth, proactive, aligned",
        "⚠ Concern: Multiple clarification loops risk slowing progress",
        "   → Now controlled via grouped queries & alignment log",
        "✅ Overall: On track, foundational logic validated"
    ], "🎯")
])

# Slide 7: Executive Summary
add_content_slide(prs, "6️⃣ Executive Summary", [
    ("Key Takeaways", [
        "✅ Model functional with core scheduling logic",
        "🔧 Demand & synthetic BOM alignment in progress",
        "❗ Awaiting one Biocon clarification (changeover rules)",
        "💼 Ready for alignment with Iris & ORMAE next week",
        "📈 Expected completion: Full Q4 2025 planning by end of week"
    ], "📌")
])

# Save presentation
output_path = "/home/supriyo/Downloads/Biocon_Daily_Production_Planning_Status.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📅 Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")