from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme - Tech/AI themed
HEADER_BG = RGBColor(26, 35, 126)  # Deep Indigo
ACCENT = RGBColor(33, 150, 243)    # Blue
TEXT = RGBColor(33, 33, 33)        # Dark Gray
SUCCESS = RGBColor(76, 175, 80)    # Green
WARNING = RGBColor(255, 152, 0)    # Orange
CODE_BG = RGBColor(245, 245, 245)  # Light Gray

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = HEADER_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(187, 222, 251)

def add_content_slide(prs, title, content_items):
    """Add content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = HEADER_BG
    title_shape.line.color.rgb = HEADER_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    y_pos = 1.0
    for item in content_items:
        if isinstance(item, dict):
            item_type = item.get('type', 'text')
            
            if item_type == 'heading':
                heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
                heading_frame = heading_box.text_frame
                p = heading_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT
                y_pos += 0.45
                
            elif item_type == 'bullet':
                bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.7), Inches(0.28))
                bullet_frame = bullet_box.text_frame
                bullet_frame.word_wrap = True
                p = bullet_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT
                p.level = item.get('level', 0)
                y_pos += 0.30
                
            elif item_type == 'code':
                code_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.7), Inches(0.35))
                code_frame = code_box.text_frame
                code_frame.word_wrap = True
                p = code_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(9)
                p.font.name = 'Courier New'
                p.font.color.rgb = RGBColor(0, 100, 0)
                y_pos += 0.38

# Slide 1: Title
add_title_slide(prs, 
    "🔍 RAG System Observability\n& Evaluation", 
    "Continuous Validation with CI/CD Integration\nUsing DeepEval, RAGAS & Kerb")

# Slide 2: Overview
add_content_slide(prs, "📊 Observability & Evaluation Framework", [
    {'type': 'heading', 'text': '🎯 Key Objectives'},
    {'type': 'bullet', 'text': 'Ensure RAG system quality through automated evaluation'},
    {'type': 'bullet', 'text': 'Monitor performance with SLIs/SLOs (Service Level Indicators/Objectives)'},
    {'type': 'bullet', 'text': 'Integrate quality gates into CI/CD pipelines'},
    {'type': 'bullet', 'text': 'Track cost, latency, and response quality metrics'},
    
    {'type': 'heading', 'text': '🛠️ Technology Stack'},
    {'type': 'bullet', 'text': 'DeepEval: Comprehensive LLM evaluation framework'},
    {'type': 'bullet', 'text': 'RAGAS Metrics: Specialized RAG evaluation (precision, recall, faithfulness)'},
    {'type': 'bullet', 'text': 'Kerb: CI/CD orchestration for ML/AI pipelines'},
    {'type': 'bullet', 'text': 'Traces: End-to-end request tracking and debugging'},
])

# Slide 3: RAGAS Metrics with DeepEval
add_content_slide(prs, "📏 RAGAS Metrics Implementation", [
    {'type': 'heading', 'text': '🔧 Core Metrics from DeepEval'},
    {'type': 'bullet', 'text': '✅ Contextual Precision: Relevance of retrieved context to query'},
    {'type': 'bullet', 'text': '✅ Contextual Recall: Coverage of relevant information in context'},
    {'type': 'bullet', 'text': '✅ Answer Relevancy: How well answer addresses the question'},
    {'type': 'bullet', 'text': '✅ Faithfulness: Answer grounding in retrieved context (no hallucinations)'},
    {'type': 'bullet', 'text': '✅ Bias Detection: Identify potential biases in responses'},
    
    {'type': 'heading', 'text': '💻 Code Example'},
    {'type': 'code', 'text': 'from deepeval.metrics.ragas import ('},
    {'type': 'code', 'text': '    RAGASContextualPrecisionMetric,'},
    {'type': 'code', 'text': '    RAGASContextualRecallMetric,'},
    {'type': 'code', 'text': '    RAGASAnswerRelevancyMetric,'},
    {'type': 'code', 'text': '    RAGASFaithfulnessMetric'},
    {'type': 'code', 'text': ')'},
    {'type': 'code', 'text': 'from deepeval.metrics import BiasMetric'},
])

# Slide 4: Test Case Structure
add_content_slide(prs, "🧪 LLM Test Case Implementation", [
    {'type': 'heading', 'text': '📝 Test Case Structure'},
    {'type': 'code', 'text': 'from deepeval.test_case import LLMTestCase'},
    {'type': 'code', 'text': 'from deepeval import evaluate'},
    {'type': 'code', 'text': ''},
    {'type': 'code', 'text': 'test_case = LLMTestCase('},
    {'type': 'code', 'text': '    input="What is retrieval-augmented generation?",'},
    {'type': 'code', 'text': '    actual_output=model_response,'},
    {'type': 'code', 'text': '    expected_output=ground_truth,'},
    {'type': 'code', 'text': '    context=retrieved_documents,'},
    {'type': 'code', 'text': '    retrieval_context=all_available_docs'},
    {'type': 'code', 'text': ')'},
    
    {'type': 'heading', 'text': '⚡ Metric Execution'},
    {'type': 'code', 'text': 'metrics = [contextual_precision, contextual_recall,'},
    {'type': 'code', 'text': '          answer_relevancy, faithfulness, bias]'},
    {'type': 'code', 'text': 'evaluate(test_cases=[test_case], metrics=metrics)'},
])

# Slide 5: CI/CD Integration
add_content_slide(prs, "🔄 CI/CD Pipeline Integration", [
    {'type': 'heading', 'text': '🚀 Continuous Validation Workflow'},
    {'type': 'bullet', 'text': '1. Code Commit → Trigger automated test suite'},
    {'type': 'bullet', 'text': '2. Run DeepEval metrics on test dataset'},
    {'type': 'bullet', 'text': '3. Check against SLO thresholds (e.g., Faithfulness > 0.85)'},
    {'type': 'bullet', 'text': '4. Quality Gate: Pass/Fail based on metrics'},
    {'type': 'bullet', 'text': '5. Deploy only if all gates pass'},
    
    {'type': 'heading', 'text': '📊 Example SLO Thresholds'},
    {'type': 'bullet', 'text': '• Contextual Precision ≥ 0.80'},
    {'type': 'bullet', 'text': '• Faithfulness ≥ 0.85 (critical: no hallucinations)'},
    {'type': 'bullet', 'text': '• Answer Relevancy ≥ 0.75'},
    {'type': 'bullet', 'text': '• Latency (p95) ≤ 2000ms'},
    {'type': 'bullet', 'text': '• Cost per request ≤ $0.05'},
])

# Slide 6: Kerb Framework
add_content_slide(prs, "⚙️ Kerb: CI/CD for RAG Systems", [
    {'type': 'heading', 'text': '🏗️ Kerb Framework Benefits'},
    {'type': 'bullet', 'text': 'Specialized for Large-Scale Distributed RAG systems'},
    {'type': 'bullet', 'text': 'Orchestration of complex evaluation pipelines'},
    {'type': 'bullet', 'text': 'Integration with monitoring tools (Prometheus, Grafana)'},
    {'type': 'bullet', 'text': 'Automated rollback on metric degradation'},
    
    {'type': 'heading', 'text': '📚 Reference Architecture'},
    {'type': 'bullet', 'text': 'Source: APXML Large-Scale Distributed RAG Course'},
    {'type': 'bullet', 'text': 'Chapter 5: Orchestration & Operationalization'},
    {'type': 'bullet', 'text': 'Topic: CI/CD Pipelines for RAG Systems'},
    {'type': 'bullet', 'text': 'URL: apxml.com/courses/large-scale-distributed-rag/'},
    {'type': 'bullet', 'text': '      chapter-5-orchestration-operationalization-large-scale-rag/'},
    {'type': 'bullet', 'text': '      ci-cd-pipelines-rag-systems'},
])

# Slide 7: Tracing & Observability
add_content_slide(prs, "🔍 Traces & Real-Time Monitoring", [
    {'type': 'heading', 'text': '📈 Observability Components'},
    {'type': 'bullet', 'text': 'End-to-End Traces: Track request flow from query to response'},
    {'type': 'bullet', 'text': 'Span Analysis: Measure retrieval, LLM inference, post-processing times'},
    {'type': 'bullet', 'text': 'Cost Tracking: Monitor token usage and API costs per request'},
    {'type': 'bullet', 'text': 'Error Logging: Capture failures, timeouts, and edge cases'},
    
    {'type': 'heading', 'text': '⚡ Key Metrics Dashboard'},
    {'type': 'bullet', 'text': '• Request Rate (req/sec)'},
    {'type': 'bullet', 'text': '• Average Latency (p50, p95, p99)'},
    {'type': 'bullet', 'text': '• Quality Score Trends (Faithfulness over time)'},
    {'type': 'bullet', 'text': '• Cost per 1K requests'},
    {'type': 'bullet', 'text': '• Error Rate & Types'},
])

# Slide 8: Complete Implementation Example
add_content_slide(prs, "💡 End-to-End Implementation", [
    {'type': 'heading', 'text': '📦 Complete Evaluation Pipeline'},
    {'type': 'code', 'text': '# 1. Initialize Metrics'},
    {'type': 'code', 'text': 'contextual_precision = RAGASContextualPrecisionMetric()'},
    {'type': 'code', 'text': 'contextual_recall = RAGASContextualRecallMetric()'},
    {'type': 'code', 'text': 'answer_relevancy = RAGASAnswerRelevancyMetric()'},
    {'type': 'code', 'text': 'faithfulness = RAGASFaithfulnessMetric()'},
    {'type': 'code', 'text': 'bias = BiasMetric()'},
    {'type': 'code', 'text': ''},
    {'type': 'code', 'text': '# 2. Create Test Cases'},
    {'type': 'code', 'text': 'test_cases = [LLMTestCase(...) for query in test_queries]'},
    {'type': 'code', 'text': ''},
    {'type': 'code', 'text': '# 3. Run Evaluation'},
    {'type': 'code', 'text': 'results = evaluate('},
    {'type': 'code', 'text': '    test_cases=test_cases,'},
    {'type': 'code', 'text': '    metrics=[contextual_precision, contextual_recall,'},
    {'type': 'code', 'text': '             answer_relevancy, faithfulness, bias]'},
    {'type': 'code', 'text': ')'},
])

# Slide 9: Best Practices
add_content_slide(prs, "✅ Best Practices & Recommendations", [
    {'type': 'heading', 'text': '🎯 Evaluation Strategy'},
    {'type': 'bullet', 'text': 'Test Dataset: Maintain diverse, representative test set (100+ examples)'},
    {'type': 'bullet', 'text': 'Regular Updates: Refresh test cases with production edge cases'},
    {'type': 'bullet', 'text': 'Balanced Metrics: Don\'t optimize single metric at expense of others'},
    {'type': 'bullet', 'text': 'Human-in-Loop: Periodically validate automated scores'},
    
    {'type': 'heading', 'text': '⚠️ Common Pitfalls to Avoid'},
    {'type': 'bullet', 'text': '❌ Overfitting to test set'},
    {'type': 'bullet', 'text': '❌ Ignoring cost/latency trade-offs'},
    {'type': 'bullet', 'text': '❌ Not monitoring production metrics post-deployment'},
    {'type': 'bullet', 'text': '❌ Setting overly strict SLOs that block valid improvements'},
])

# Slide 10: Summary & Next Steps
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide10.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 248, 255)

title_box = slide10.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "🎯 Summary & Implementation Roadmap"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = HEADER_BG

content_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

items = [
    ("✅ KEY TAKEAWAYS", [
        "DeepEval + RAGAS provide comprehensive RAG evaluation",
        "5 critical metrics: Precision, Recall, Relevancy, Faithfulness, Bias",
        "CI/CD integration ensures continuous quality validation",
        "SLIs/SLOs enable data-driven performance gates"
    ]),
    ("🔧 IMPLEMENTATION STEPS", [
        "1. Set up DeepEval with RAGAS metrics",
        "2. Create representative test dataset",
        "3. Define SLO thresholds for each metric",
        "4. Integrate with CI/CD (Kerb framework)",
        "5. Set up dashboards for real-time monitoring"
    ]),
    ("📚 RESOURCES", [
        "DeepEval Docs: docs.confident-ai.com",
        "RAGAS Framework: docs.ragas.io",
        "APXML RAG Course: apxml.com/courses/large-scale-distributed-rag",
        "Kerb Documentation: [Framework-specific docs]"
    ])
]

for section_title, section_items in items:
    p = content_frame.add_paragraph()
    p.text = section_title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.level = 0
    
    for item in section_items:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.level = 1

# Save
output_path = "/home/supriyo/repogen/RAG_Observability_Evaluation_Framework.pptx"
prs.save(output_path)
print(f"✅ PowerPoint created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"\n📌 Slide Contents:")
print(f"   1. Title: RAG System Observability & Evaluation")
print(f"   2. Overview: Framework objectives and tech stack")
print(f"   3. RAGAS Metrics: Implementation with DeepEval")
print(f"   4. Test Case Structure: LLMTestCase examples")
print(f"   5. CI/CD Integration: Continuous validation workflow")
print(f"   6. Kerb Framework: CI/CD orchestration for RAG")
print(f"   7. Traces & Monitoring: Real-time observability")
print(f"   8. Implementation Example: End-to-end code")
print(f"   9. Best Practices: Strategy and pitfalls")
print(f"   10. Summary: Takeaways and roadmap")